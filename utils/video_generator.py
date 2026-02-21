#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
视频生成器模块
基于PPT内容生成有声视频

功能：
1. 读取PPT内容或直接从文本生成PPT
2. 使用PDF中转方案（poppler-utils）或Pillow生成幻灯片图片
3. 使用Edge-TTS生成语音
4. 使用ffmpeg合成最终视频
"""

import os
import sys
import json
import asyncio
import argparse
import tempfile
import subprocess
from datetime import datetime
from pathlib import Path

# 添加项目根目录到路径
script_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(script_dir)
if project_root not in sys.path:
    sys.path.insert(0, project_root)

try:
    import edge_tts
except ImportError:
    print("⚠️ 请安装 edge-tts: pip install edge-tts")
    edge_tts = None

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("⚠️ 请安装 Pillow: pip install Pillow")
    Image = None

# 导入PPT生成器
try:
    from utils.ppt_generator import generate_ppt, get_style_config
except ImportError:
    print("⚠️ 无法导入PPT生成器")
    generate_ppt = None


# ============================================================================
# 配置
# ============================================================================

# 视频配置
VIDEO_CONFIG = {
    'width': 1920,
    'height': 1080,
    'fps': 24,
    'default_duration': 5,  # 默认每页时长（秒）
    'use_pdf_method': True,  # 优先使用PDF中转方案
}

# 检查poppler-utils是否可用
def check_poppler_available():
    """检查poppler-utils是否可用"""
    try:
        result = subprocess.run(
            ['pdftoppm', '-v'],
            capture_output=True,
            timeout=5
        )
        return result.returncode == 0
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False

# 检查reportlab是否可用（用于生成PDF）
def check_reportlab_available():
    """检查reportlab是否可用"""
    try:
        import reportlab
        return True
    except ImportError:
        return False


def letterbox_resize(img, target_size, fill_color=(255, 255, 255)):
    """保持比例调整图片大小，使用letterbox方式填充
    
    Args:
        img: PIL Image对象
        target_size: 目标尺寸 (width, height)
        fill_color: 填充颜色
    
    Returns:
        PIL Image对象
    """
    target_width, target_height = target_size
    original_width, original_height = img.size
    
    # 计算缩放比例
    scale = min(target_width / original_width, target_height / original_height)
    new_width = int(original_width * scale)
    new_height = int(original_height * scale)
    
    # 缩放图片
    img_resized = img.resize((new_width, new_height), Image.LANCZOS)
    
    # 创建目标尺寸的背景
    result = Image.new('RGB', target_size, fill_color)
    
    # 居中放置缩放后的图片
    offset_x = (target_width - new_width) // 2
    offset_y = (target_height - new_height) // 2
    result.paste(img_resized, (offset_x, offset_y))
    
    return result

# Edge-TTS中文语音选项
TTS_VOICES = {
    'female': 'zh-CN-XiaoxiaoNeural',        # 女声 - 晓晓
    'male': 'zh-CN-YunxiNeural',              # 男声 - 云希
    'female_news': 'zh-CN-XiaoyiNeural',      # 女声新闻 - 晓伊
    'male_news': 'zh-CN-YunjianNeural',       # 男声新闻 - 云健
    'female_gentle': 'zh-CN-XiaochenNeural',  # 女声温柔 - 晓辰
    'male_cheerful': 'zh-CN-YunyangNeural',   # 男声活泼 - 云扬
}

# 默认语音
DEFAULT_VOICE = 'female'


# ============================================================================
# PDF生成功能（PDF中转方案）
# ============================================================================

def generate_pdf_from_slides(slides_data, output_path, style='business'):
    """从幻灯片数据生成PDF文件
    
    Args:
        slides_data (list): 幻灯片数据列表
        output_path (str): 输出PDF路径
        style (str): 风格
    
    Returns:
        str: PDF文件路径
    """
    try:
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.units import cm, inch
        from reportlab.lib.utils import simpleSplit
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except ImportError:
        raise RuntimeError("请安装 reportlab: pip install reportlab")
    
    style_config = get_style_config(style)
    
    # 创建PDF文档（使用16:9尺寸以匹配视频比例）
    # 1920x1080 比例转换为英寸：约 16" x 9"
    page_width = 16 * inch
    page_height = 9 * inch
    pagesize = (page_width, page_height)
    
    doc = SimpleDocTemplate(
        output_path,
        pagesize=pagesize,
        rightMargin=0.5*inch,
        leftMargin=0.5*inch,
        topMargin=0.5*inch,
        bottomMargin=0.5*inch
    )
    
    # 注册中文字体
    font_paths = [
        '/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc',
        '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
        '/System/Library/Fonts/PingFang.ttc',
        'C:/Windows/Fonts/msyh.ttc',
    ]
    
    font_name = 'CustomChinese'
    font_registered = False
    
    for font_path in font_paths:
        if os.path.exists(font_path):
            try:
                pdfmetrics.registerFont(TTFont(font_name, font_path))
                font_registered = True
                break
            except:
                continue
    
    if not font_registered:
        print("⚠️ 未找到中文字体，使用默认字体")
        font_name = 'Helvetica'
    
    # 样式配置
    styles = getSampleStyleSheet()
    
    # 转换RGB颜色
    title_color = colors.HexColor('#{:02x}{:02x}{:02x}'.format(
        style_config['title_color'][0],
        style_config['title_color'][1],
        style_config['title_color'][2]
    ))
    
    text_color = colors.HexColor('#{:02x}{:02x}{:02x}'.format(
        style_config['text_color'][0],
        style_config['text_color'][1],
        style_config['text_color'][2]
    ))
    
    accent_color = colors.HexColor('#{:02x}{:02x}{:02x}'.format(
        style_config['accent_color'][0],
        style_config['accent_color'][1],
        style_config['accent_color'][2]
    ))
    
    # 自定义样式
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontName=font_name,
        fontSize=28,
        textColor=title_color,
        spaceAfter=20,
        alignment=1  # 居中
    )
    
    subtitle_style = ParagraphStyle(
        'CustomSubtitle',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=18,
        textColor=text_color,
        spaceAfter=30,
        alignment=1  # 居中
    )
    
    content_title_style = ParagraphStyle(
        'ContentTitle',
        parent=styles['Heading2'],
        fontName=font_name,
        fontSize=24,
        textColor=title_color,
        spaceAfter=15
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=16,
        textColor=text_color,
        spaceAfter=12,
        leftIndent=20,
        bulletIndent=10
    )
    
    date_style = ParagraphStyle(
        'CustomDate',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=12,
        textColor=accent_color,
        spaceBefore=50,
        alignment=1  # 居中
    )
    
    conclusion_style = ParagraphStyle(
        'CustomConclusion',
        parent=styles['Heading2'],
        fontName=font_name,
        fontSize=24,
        textColor=title_color,
        spaceAfter=20,
        spaceBefore=30,
        alignment=1  # 居中
    )
    
    thanks_style = ParagraphStyle(
        'CustomThanks',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=20,
        textColor=accent_color,
        spaceBefore=40,
        alignment=1  # 居中
    )
    
    # 构建PDF内容
    story = []
    
    for slide in slides_data:
        slide_type = slide.get('type', 'content')
        
        if slide_type == 'title':
            # 封面页
            story.append(Spacer(1, 2*cm))
            story.append(Paragraph(slide.get('title', ''), title_style))
            
            if slide.get('subtitle'):
                story.append(Paragraph(slide.get('subtitle', ''), subtitle_style))
            
            # 添加日期
            date_text = datetime.now().strftime('%Y年%m月%d日')
            story.append(Paragraph(date_text, date_style))
            story.append(Spacer(1, 1*cm))
            
            # 添加分隔线
            story.append(Table([['']], colWidths=[16*cm], style=[
                ('LINEBELOW', (0, 0), (-1, -1), 2, accent_color)
            ]))
            
        elif slide_type == 'conclusion':
            # 结尾页
            story.append(Spacer(1, 2*cm))
            story.append(Paragraph('总结', conclusion_style))
            
            if slide.get('conclusion'):
                story.append(Paragraph(slide.get('conclusion', ''), body_style))
            
            story.append(Paragraph('感谢观看', thanks_style))
            
        else:
            # 内容页
            story.append(Paragraph(slide.get('title', ''), content_title_style))
            
            for point in slide.get('points', []):
                story.append(Paragraph(f"• {point}", body_style))
        
        # 添加分页符（除了最后一页）
        if slide != slides_data[-1]:
            from reportlab.platypus import PageBreak
            story.append(PageBreak())
    
    # 生成PDF
    doc.build(story)
    
    return output_path


def convert_pdf_to_images(pdf_path, output_dir, dpi=300):
    """将PDF转换为图片序列（使用pdftoppm）
    
    Args:
        pdf_path (str): PDF文件路径
        output_dir (str): 输出目录
        dpi (int): DPI值（默认300）
    
    Returns:
        list: 图片文件路径列表
    """
    if not check_poppler_available():
        raise RuntimeError("poppler-utils未安装，请运行: sudo apt-get install poppler-utils")
    
    os.makedirs(output_dir, exist_ok=True)
    
    # 使用pdftoppm转换PDF为PNG图片
    base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    
    cmd = [
        'pdftoppm',
        '-png',
        '-r', str(dpi),
        pdf_path,
        os.path.join(output_dir, base_name)
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    
    if result.returncode != 0:
        raise RuntimeError(f"PDF转换失败: {result.stderr}")
    
    # 收集生成的图片文件
    image_files = []
    for file in sorted(os.listdir(output_dir)):
        if file.startswith(base_name) and file.endswith('.png'):
            image_files.append(os.path.join(output_dir, file))
    
    if not image_files:
        raise RuntimeError("未生成任何图片文件")
    
    return image_files


# ============================================================================
# 幻灯片图片生成（Pillow备选方案）
# ============================================================================

def create_slide_image(title, points, style='business', output_path=None, index=0):
    """创建幻灯片图片
    
    Args:
        title (str): 标题
        points (list): 要点列表
        style (str): 风格
        output_path (str): 输出路径
        index (int): 幻灯片索引
    
    Returns:
        str: 图片路径
    """
    if Image is None:
        raise RuntimeError("Pillow库未安装")
    
    # 获取样式配置
    style_config = get_style_config(style)
    
    # 创建图片
    width, height = VIDEO_CONFIG['width'], VIDEO_CONFIG['height']
    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    # 获取渐变背景颜色
    gradient_config = style_config.get('gradient', {})
    if gradient_config.get('type') == 'linear':
        stops = gradient_config.get('stops', [])
        if len(stops) >= 2:
            # 使用中间色作为背景
            bg_color = stops[len(stops)//2][1]
            img = Image.new('RGB', (width, height), color=(
                bg_color[0], bg_color[1], bg_color[2]
            ))
            draw = ImageDraw.Draw(img)
    elif gradient_config.get('type') == 'solid':
        bg_color = gradient_config['color']
        img = Image.new('RGB', (width, height), color=(
            bg_color[0], bg_color[1], bg_color[2]
        ))
        draw = ImageDraw.Draw(img)
    
    # 字体配置
    title_color = style_config['title_color']
    text_color = style_config['text_color']
    accent_color = style_config['accent_color']
    
    # 尝试加载字体
    try:
        title_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 56)
        body_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 32)
        small_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 24)
    except:
        try:
            title_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc', 56)
            body_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc', 32)
            small_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc', 24)
        except:
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
            small_font = ImageFont.load_default()
    
    # 添加顶部装饰条
    decorations = style_config.get('decorations', {})
    if 'side_bar' in decorations and decorations['side_bar'].get('enabled'):
        bar_color = decorations['side_bar']['color']
        bar_width = int(width * decorations['side_bar']['width'])
        draw.rectangle([width - bar_width, 0, width, height], fill=(
            bar_color[0], bar_color[1], bar_color[2]
        ))
    
    if 'side_glow' in decorations and decorations['side_glow'].get('enabled'):
        glow_color = decorations['side_glow']['color']
        glow_width = int(width * decorations['side_glow']['width'])
        draw.rectangle([0, 0, glow_width, height], fill=(
            glow_color[0], glow_color[1], glow_color[2]
        ))
    
    # 绘制标题
    title_y = 100
    draw.text((100, title_y), title, font=title_font, fill=(
        title_color[0], title_color[1], title_color[2]
    ))
    
    # 绘制分隔线
    line_y = title_y + 80
    draw.rectangle([100, line_y, width - 200, line_y + 4], fill=(
        accent_color[0], accent_color[1], accent_color[2]
    ))
    
    # 绘制要点
    point_y = line_y + 60
    for i, point in enumerate(points[:6]):  # 最多6个要点
        text = f"• {point}"
        draw.text((100, point_y), text, font=body_font, fill=(
            text_color[0], text_color[1], text_color[2]
        ))
        point_y += 60
    
    # 绘制页码
    page_text = f"{index + 1}"
    draw.text((width - 150, height - 80), page_text, font=small_font, fill=(
        accent_color[0], accent_color[1], accent_color[2]
    ))
    
    # 添加底部装饰线
    if 'footer_line' in decorations and decorations['footer_line'].get('enabled'):
        footer_color = decorations['footer_line']['color']
        footer_height = int(height * decorations['footer_line']['height'])
        draw.rectangle([0, height - footer_height - 30, width, height], fill=(
            footer_color[0], footer_color[1], footer_color[2]
        ))
    
    # 保存图片
    if output_path is None:
        output_dir = os.path.join(script_dir, '..', 'output', 'video_temp')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'slide_{index:03d}.png')
    
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    img.save(output_path, 'PNG')
    
    return output_path


def create_title_slide_image(title, subtitle, style='business', output_path=None):
    """创建封面页图片
    
    Args:
        title (str): 主标题
        subtitle (str): 副标题
        style (str): 风格
        output_path (str): 输出路径
    
    Returns:
        str: 图片路径
    """
    if Image is None:
        raise RuntimeError("Pillow库未安装")
    
    style_config = get_style_config(style)
    width, height = VIDEO_CONFIG['width'], VIDEO_CONFIG['height']
    
    # 创建背景
    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    # 设置背景颜色
    gradient_config = style_config.get('gradient', {})
    if gradient_config.get('type') == 'linear':
        stops = gradient_config.get('stops', [])
        if len(stops) >= 2:
            bg_color = stops[len(stops)//2][1]
            img = Image.new('RGB', (width, height), color=(
                bg_color[0], bg_color[1], bg_color[2]
            ))
            draw = ImageDraw.Draw(img)
    elif gradient_config.get('type') == 'solid':
        bg_color = gradient_config['color']
        img = Image.new('RGB', (width, height), color=(
            bg_color[0], bg_color[1], bg_color[2]
        ))
        draw = ImageDraw.Draw(img)
    
    title_color = style_config['title_color']
    text_color = style_config['text_color']
    accent_color = style_config['accent_color']
    
    # 字体
    try:
        title_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 72)
        subtitle_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 36)
        date_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 24)
    except:
        try:
            title_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc', 72)
            subtitle_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc', 36)
            date_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc', 24)
        except:
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()
            date_font = ImageFont.load_default()
    
    # 装饰
    decorations = style_config.get('decorations', {})
    if 'side_bar' in decorations and decorations['side_bar'].get('enabled'):
        bar_color = decorations['side_bar']['color']
        bar_width = int(width * decorations['side_bar']['width'])
        draw.rectangle([width - bar_width, 0, width, height], fill=(
            bar_color[0], bar_color[1], bar_color[2]
        ))
    
    # 绘制标题（居中）
    title_y = height // 2 - 100
    draw.text((width // 2, title_y), title, font=title_font, fill=(
        title_color[0], title_color[1], title_color[2]
    ), anchor='mm')
    
    # 绘制副标题
    if subtitle:
        subtitle_y = title_y + 80
        draw.text((width // 2, subtitle_y), subtitle, font=subtitle_font, fill=(
            text_color[0], text_color[1], text_color[2]
        ), anchor='mm')
    
    # 绘制日期
    date_text = datetime.now().strftime('%Y年%m月%d日')
    date_y = height - 150
    draw.text((width // 2, date_y), date_text, font=date_font, fill=(
        accent_color[0], accent_color[1], accent_color[2]
    ), anchor='mm')
    
    # 底部装饰线
    if 'footer_line' in decorations and decorations['footer_line'].get('enabled'):
        footer_color = decorations['footer_line']['color']
        footer_height = int(height * decorations['footer_line']['height'])
        draw.rectangle([0, height - footer_height - 30, width, height], fill=(
            footer_color[0], footer_color[1], footer_color[2]
        ))
    
    if output_path is None:
        output_dir = os.path.join(script_dir, '..', 'output', 'video_temp')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, 'slide_000.png')
    
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    img.save(output_path, 'PNG')
    
    return output_path


def create_conclusion_slide_image(conclusion, style='business', output_path=None, index=0):
    """创建结尾页图片
    
    Args:
        conclusion (str): 总结内容
        style (str): 风格
        output_path (str): 输出路径
        index (int): 幻灯片索引
    
    Returns:
        str: 图片路径
    """
    if Image is None:
        raise RuntimeError("Pillow库未安装")
    
    style_config = get_style_config(style)
    width, height = VIDEO_CONFIG['width'], VIDEO_CONFIG['height']
    
    # 创建背景
    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    # 设置背景颜色
    gradient_config = style_config.get('gradient', {})
    if gradient_config.get('type') == 'linear':
        stops = gradient_config.get('stops', [])
        if len(stops) >= 2:
            bg_color = stops[len(stops)//2][1]
            img = Image.new('RGB', (width, height), color=(
                bg_color[0], bg_color[1], bg_color[2]
            ))
            draw = ImageDraw.Draw(img)
    elif gradient_config.get('type') == 'solid':
        bg_color = gradient_config['color']
        img = Image.new('RGB', (width, height), color=(
            bg_color[0], bg_color[1], bg_color[2]
        ))
        draw = ImageDraw.Draw(img)
    
    title_color = style_config['title_color']
    text_color = style_config['text_color']
    accent_color = style_config['accent_color']
    
    # 字体
    try:
        title_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 56)
        body_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 36)
        thanks_font = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 48)
    except:
        try:
            title_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc', 56)
            body_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc', 36)
            thanks_font = ImageFont.truetype('/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc', 48)
        except:
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
            thanks_font = ImageFont.load_default()
    
    # 装饰
    decorations = style_config.get('decorations', {})
    if 'side_bar' in decorations and decorations['side_bar'].get('enabled'):
        bar_color = decorations['side_bar']['color']
        bar_width = int(width * decorations['side_bar']['width'])
        draw.rectangle([width - bar_width, 0, width, height], fill=(
            bar_color[0], bar_color[1], bar_color[2]
        ))
    
    # 绘制"总结"标题
    title_y = height // 2 - 150
    draw.text((width // 2, title_y), '总结', font=title_font, fill=(
        title_color[0], title_color[1], title_color[2]
    ), anchor='mm')
    
    # 绘制总结内容
    if conclusion:
        conclusion_y = title_y + 100
        draw.text((width // 2, conclusion_y), conclusion, font=body_font, fill=(
            text_color[0], text_color[1], text_color[2]
        ), anchor='mm')
    
    # 绘制感谢语
    thanks_y = height - 200
    draw.text((width // 2, thanks_y), '感谢观看', font=thanks_font, fill=(
        accent_color[0], accent_color[1], accent_color[2]
    ), anchor='mm')
    
    if output_path is None:
        output_dir = os.path.join(script_dir, '..', 'output', 'video_temp')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'slide_{index:03d}.png')
    
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    img.save(output_path, 'PNG')
    
    return output_path


# ============================================================================
# TTS语音生成
# ============================================================================

async def generate_audio_async(text, output_path, voice='female'):
    """异步生成语音文件
    
    Args:
        text (str): 要转换的文本
        output_path (str): 输出路径
        voice (str): 语音类型
    """
    if edge_tts is None:
        raise RuntimeError("edge-tts库未安装")
    
    voice_name = TTS_VOICES.get(voice, TTS_VOICES[DEFAULT_VOICE])
    
    communicate = edge_tts.Communicate(text, voice_name)
    await communicate.save(output_path)


def generate_audio(text, output_path, voice='female'):
    """生成语音文件（同步接口）
    
    Args:
        text (str): 要转换的文本
        output_path (str): 输出路径
        voice (str): 语音类型
    
    Returns:
        str: 音频文件路径
    """
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    asyncio.run(generate_audio_async(text, output_path, voice))
    return output_path


def get_audio_duration(audio_path):
    """获取音频时长
    
    Args:
        audio_path (str): 音频文件路径
    
    Returns:
        float: 时长（秒）
    """
    try:
        result = subprocess.run(
            ['ffprobe', '-v', 'error', '-show_entries', 'format=duration',
             '-of', 'default=noprint_wrappers=1:nokey=1', audio_path],
            capture_output=True, text=True, check=True
        )
        return float(result.stdout.strip())
    except:
        return VIDEO_CONFIG['default_duration']


# ============================================================================
# 视频合成
# ============================================================================

def create_video_from_slides(slides_data, output_path, style='business', voice='female', existing_images=None):
    """从幻灯片数据创建视频
    
    Args:
        slides_data (list): 幻灯片数据列表，每项包含:
            - type: 'title', 'content', 'conclusion'
            - title: 标题
            - subtitle/content/conclusion: 内容
            - points: 要点列表
        output_path (str): 输出视频路径
        style (str): 风格
        voice (str): 语音类型
        existing_images (list): 已存在的图片路径列表（可选）
    
    Returns:
        str: 视频文件路径
    """
    if not os.path.exists('/usr/bin/ffmpeg'):
        raise RuntimeError("ffmpeg未安装")
    
    # 创建临时目录
    temp_dir = tempfile.mkdtemp(prefix='video_')
    
    try:
        # 检查是否使用PDF中转方案
        use_pdf = VIDEO_CONFIG.get('use_pdf_method', True)
        use_pdf = use_pdf and check_poppler_available() and check_reportlab_available()
        
        if use_pdf:
            print("  📄 使用PDF中转方案生成幻灯片")
        else:
            print("  🎨 使用Pillow方案生成幻灯片")
            if not check_poppler_available():
                print("    ⚠️ poppler-utils未安装，将使用Pillow备选方案")
            if not check_reportlab_available():
                print("    ⚠️ reportlab未安装，将使用Pillow备选方案")
        
        # 生成图片和音频
        slide_files = []
        audio_files = []
        durations = []
        
        if existing_images:
            # 使用已存在的图片（来自PPTX转换）
            slide_files = existing_images
        elif use_pdf:
            pdf_path = os.path.join(temp_dir, 'slides.pdf')
            generate_pdf_from_slides(slides_data, pdf_path, style)
            
            # 2. 转换PDF为图片序列
            images_temp_dir = os.path.join(temp_dir, 'pdf_images')
            slide_files = convert_pdf_to_images(pdf_path, images_temp_dir, dpi=150)
            
            # 确保图片数量与幻灯片数量一致
            if len(slide_files) != len(slides_data):
                print(f"  ⚠️ PDF生成的图片数量({len(slide_files)})与幻灯片数量({len(slides_data)})不一致")
                # 使用Pillow方案作为备选
                use_pdf = False
            else:
                # 调整图片尺寸为1920x1080（使用letterbox方式保持比例）
                if Image:
                    for i, img_path in enumerate(slide_files):
                        img = Image.open(img_path)
                        img_resized = letterbox_resize(img, (1920, 1080))
                        img_resized.save(img_path, 'PNG')
        
        if not use_pdf:
            # 使用Pillow方案（备选）
            for i, slide in enumerate(slides_data):
                slide_type = slide.get('type', 'content')
                
                # 生成图片
                if slide_type == 'title':
                    img_path = os.path.join(temp_dir, f'slide_{i:03d}.png')
                    create_title_slide_image(
                        slide.get('title', ''),
                        slide.get('subtitle', ''),
                        style=style,
                        output_path=img_path
                    )
                    # 生成音频
                    audio_text = f"{slide.get('title', '')}。{slide.get('subtitle', '')}"
                elif slide_type == 'conclusion':
                    img_path = os.path.join(temp_dir, f'slide_{i:03d}.png')
                    create_conclusion_slide_image(
                        slide.get('conclusion', ''),
                        style=style,
                        output_path=img_path,
                        index=i
                    )
                    # 生成音频
                    audio_text = f"总结。{slide.get('conclusion', '')}。感谢观看。"
                else:
                    img_path = os.path.join(temp_dir, f'slide_{i:03d}.png')
                    create_slide_image(
                        slide.get('title', ''),
                        slide.get('points', []),
                        style=style,
                        output_path=img_path,
                        index=i
                    )
                    # 生成音频
                    audio_text = slide.get('title', '')
                    for point in slide.get('points', []):
                        audio_text += f"。{point}"
                
                slide_files.append(img_path)
        
        # 生成语音
        for i, slide in enumerate(slides_data):
            slide_type = slide.get('type', 'content')
            
            if slide_type == 'title':
                audio_text = f"{slide.get('title', '')}。{slide.get('subtitle', '')}"
            elif slide_type == 'conclusion':
                audio_text = f"总结。{slide.get('conclusion', '')}。感谢观看。"
            else:
                audio_text = slide.get('title', '')
                for point in slide.get('points', []):
                    audio_text += f"。{point}"
            
            # 生成语音
            if audio_text.strip():
                audio_path = os.path.join(temp_dir, f'audio_{i:03d}.mp3')
                generate_audio(audio_text, audio_path, voice)
                audio_files.append(audio_path)
                duration = get_audio_duration(audio_path)
                durations.append(duration + 0.5)  # 额外0.5秒停顿
            else:
                audio_files.append(None)
                durations.append(VIDEO_CONFIG['default_duration'])
        
        # 创建视频片段
        video_parts = []
        for i, (slide, audio, duration) in enumerate(zip(slide_files, audio_files, durations)):
            part_path = os.path.join(temp_dir, f'part_{i:03d}.mp4')
            
            if audio:
                # 有音频，根据音频时长生成视频
                cmd = [
                    'ffmpeg', '-y',
                    '-loop', '1',
                    '-i', slide,
                    '-i', audio,
                    '-c:v', 'libx264',
                    '-tune', 'stillimage',
                    '-c:a', 'aac',
                    '-b:a', '192k',
                    '-pix_fmt', 'yuv420p',
                    '-shortest',
                    '-t', str(duration),
                    part_path
                ]
            else:
                # 无音频，生成静音视频
                cmd = [
                    'ffmpeg', '-y',
                    '-loop', '1',
                    '-i', slide,
                    '-c:v', 'libx264',
                    '-t', str(duration),
                    '-pix_fmt', 'yuv420p',
                    '-f', 'lavfi', '-i', 'anullsrc=channel_layout=stereo:sample_rate=44100',
                    '-shortest',
                    part_path
                ]
            
            subprocess.run(cmd, capture_output=True, check=True)
            video_parts.append(part_path)
        
        # 合并所有视频片段
        concat_file = os.path.join(temp_dir, 'concat.txt')
        with open(concat_file, 'w') as f:
            for part in video_parts:
                f.write(f"file '{part}'\n")
        
        # 最终合成
        subprocess.run([
            'ffmpeg', '-y',
            '-f', 'concat',
            '-safe', '0',
            '-i', concat_file,
            '-c:v', 'libx264',
            '-c:a', 'aac',
            '-movflags', '+faststart',
            output_path
        ], capture_output=True, check=True)
        
        return output_path
    
    finally:
        # 清理临时文件
        import shutil
        try:
            shutil.rmtree(temp_dir)
        except:
            pass


# ============================================================================
# 主函数
# ============================================================================

def generate_video_from_text(text, style='business', voice='female', output_path=None):
    """从文本生成视频
    
    Args:
        text (str): 输入文本
        style (str): PPT风格
        voice (str): 语音类型
        output_path (str): 输出路径
    
    Returns:
        dict: 生成结果
    """
    print("=" * 50)
    print("📹 视频生成器")
    print("=" * 50)
    print(f"📝 输入文本长度: {len(text)} 字符")
    print(f"🎨 PPT风格: {get_style_config(style)['name']}")
    print(f"🎤 语音类型: {voice}")
    print()
    
    # 1. 生成PPT大纲
    print("🔍 正在分析文本内容...")
    ppt_result = generate_ppt(
        text=text,
        style=style,
        output_path=None,
        enable_thinking=False
    )
    
    outline = ppt_result['outline']
    print(f"📋 提取标题: {outline.get('title', '无')}")
    print(f"📑 提取章节数: {len(outline.get('sections', []))}")
    print()
    
    # 2. 准备幻灯片数据
    print("🎨 正在准备幻灯片数据...")
    slides_data = []
    
    # 封面页
    slides_data.append({
        'type': 'title',
        'title': outline.get('title', '演示文稿'),
        'subtitle': outline.get('subtitle', '')
    })
    
    # 内容页
    for section in outline.get('sections', []):
        slides_data.append({
            'type': 'content',
            'title': section.get('title', ''),
            'points': section.get('points', [])
        })
    
    # 结尾页
    slides_data.append({
        'type': 'conclusion',
        'conclusion': outline.get('conclusion', '感谢观看')
    })
    
    # 3. 生成视频
    print("🎬 正在生成视频...")
    
    if output_path is None:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_dir = os.path.join(project_root, 'output', 'videos')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'video_{timestamp}.mp4')
    
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    
    create_video_from_slides(slides_data, output_path, style, voice)
    
    # 4. 输出结果
    print()
    print("=" * 50)
    print("✅ 视频生成成功！")
    print("=" * 50)
    print(f"📄 视频文件: {output_path}")
    print(f"📊 幻灯片数量: {len(slides_data)} 页")
    print(f"🎨 风格: {get_style_config(style)['name']}")
    print(f"🎤 语音: {voice}")
    print()
    print("💡 提示: 可使用视频播放器查看")
    print("=" * 50)
    
    return {
        'success': True,
        'video_path': output_path,
        'slide_count': len(slides_data),
        'style': style,
        'voice': voice
    }


def generate_video_from_pptx(pptx_path, voice='female', output_path=None):
    """从PPT文件生成视频（尝试多种方案）
    
    优先级:
    1. unoconv (如果可用) - 通过 LibreOffice 转换，保留完整样式
    2. 直接转换 PDF - 使用 pdf2image 转换
    
    Args:
        pptx_path (str): PPT文件路径
        voice (str): 语音类型
        output_path (str): 输出路径
    
    Returns:
        dict: 生成结果
    """
    print("=" * 50)
    print("📹 视频生成器")
    print("=" * 50)
    print(f"📄 PPT文件: {pptx_path}")
    print(f"🎤 语音类型: {voice}")
    print()
    
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("请安装 python-pptx: pip install python-pptx")
    
    # 读取PPT文件
    print("🔍 正在读取PPT文件...")
    prs = Presentation(pptx_path)
    
    print(f"📊 幻灯片数量: {len(prs.slides)} 页")
    print()
    
    # 创建临时目录
    temp_dir = tempfile.mkdtemp(prefix='video_pptx_')
    
    try:
        # 方案1: 尝试使用 unoconv
        print("🎨 尝试使用 unoconv 转换PPTX...")
        
        # 检查 unoconv 是否可用
        if subprocess.run(['which', 'unoconv'], capture_output=True).returncode == 0:
            try:
                # unoconv 会在原文件同目录生成 PDF
                pdf_dir = os.path.dirname(pptx_path) or '.'
                pdf_filename = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
                expected_pdf_path = os.path.join(pdf_dir, pdf_filename)
                
                # 先删除可能存在的旧 PDF
                if os.path.exists(expected_pdf_path):
                    os.remove(expected_pdf_path)
                
                result = subprocess.run(
                    ['unoconv', '-f', 'pdf', pptx_path],
                    capture_output=True,
                    timeout=60
                )
                
                if result.returncode == 0 and os.path.exists(expected_pdf_path):
                    print("  ✅ unoconv 转换成功")
                    
                    # 移动 PDF 到临时目录
                    pdf_path = os.path.join(temp_dir, 'slides.pdf')
                    shutil.move(expected_pdf_path, pdf_path)
                    
                    # 转换PDF为图片
                    images_temp_dir = os.path.join(temp_dir, 'pdf_images')
                    slide_files = convert_pdf_to_images(pdf_path, images_temp_dir, dpi=150)
                    
                    # 调整图片尺寸
                    if Image and slide_files:
                        for img_path in slide_files:
                            img = Image.open(img_path)
                            img_resized = letterbox_resize(img, (1920, 1080))
                            img_resized.save(img_path, 'PNG')
                    
                    # 解析幻灯片数据
                    slides_data = []
                    for i, slide in enumerate(prs.slides):
                        text_content = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                text_content.append(shape.text.strip())
                        
                        if not text_content:
                            continue
                        
                        if i == 0:
                            slides_data.append({
                                'type': 'title',
                                'title': text_content[0] if text_content else "标题",
                                'subtitle': text_content[1] if len(text_content) > 1 else ""
                            })
                        elif '总结' in ''.join(text_content) or '感谢' in ''.join(text_content):
                            slides_data.append({
                                'type': 'conclusion',
                                'conclusion': ' '.join(text_content[1:]) if len(text_content) > 1 else ""
                            })
                        else:
                            slides_data.append({
                                'type': 'content',
                                'title': text_content[0] if text_content else "",
                                'points': text_content[1:] if len(text_content) > 1 else []
                            })
                    
                    # 生成视频
                    print("🎬 正在合成视频...")
                    create_video_from_slides(slides_data, output_path, 'business', voice, 
                                             existing_images=slide_files)
                    
                    return {'success': True, 'video_path': output_path}
                else:
                    print("  ⚠️ unoconv 转换失败，使用备选方案")
            except Exception as e:
                print(f"  ⚠️ unoconv 转换出错: {e}，使用备选方案")
        else:
            print("  ⚠️ unoconv 不可用，使用备选方案")
            print("    提示: 安装 unoconv 可以保留 PPT 完整样式")
            print("    安装命令: sudo apt-get install unoconv libreoffice")
        
        # 方案2: 使用 PDF 中转方案（原有逻辑）
        # 解析幻灯片数据
        slides_data = []
        
        for i, slide in enumerate(prs.slides):
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            if not text_content:
                continue
            
            if i == 0:
                slides_data.append({
                    'type': 'title',
                    'title': text_content[0] if text_content else "标题",
                    'subtitle': text_content[1] if len(text_content) > 1 else ""
                })
            elif '总结' in ''.join(text_content) or '感谢' in ''.join(text_content):
                slides_data.append({
                    'type': 'conclusion',
                    'conclusion': ' '.join(text_content[1:]) if len(text_content) > 1 else ""
                })
            else:
                slides_data.append({
                    'type': 'content',
                    'title': text_content[0] if text_content else "",
                    'points': text_content[1:] if len(text_content) > 1 else []
                })
        
        # 生成视频
        print("🎬 正在生成视频...")
        print("  📄 使用PDF中转方案生成幻灯片")
        result = create_video_from_slides(slides_data, output_path, 'business', voice)
        
        return result
        
    finally:
        # 清理临时目录
        shutil.rmtree(temp_dir)
    
    for i, slide in enumerate(prs.slides):
        # 提取文本内容
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
        
        if not text_content:
            continue
        
        # 判断幻灯片类型
        if i == 0:
            # 封面页
            title = text_content[0] if text_content else "标题"
            subtitle = text_content[1] if len(text_content) > 1 else ""
            slides_data.append({
                'type': 'title',
                'title': title,
                'subtitle': subtitle
            })
        elif i == len(prs.slides) - 1:
            # 结尾页
            conclusion = text_content[0] if text_content else "感谢观看"
            slides_data.append({
                'type': 'conclusion',
                'conclusion': conclusion
            })
        else:
            # 内容页
            title = text_content[0] if text_content else "章节标题"
            points = text_content[1:] if len(text_content) > 1 else []
            slides_data.append({
                'type': 'content',
                'title': title,
                'points': points
            })
    
    print(f"📋 解析完成: {len(slides_data)} 张幻灯片")
    print()
    
    # 生成视频
    print("🎬 正在生成视频...")
    
    if output_path is None:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_dir = os.path.join(project_root, 'output', 'videos')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'video_{timestamp}.mp4')
    
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    
    # 使用默认商务风格
    create_video_from_slides(slides_data, output_path, 'business', voice)
    
    # 输出结果
    print()
    print("=" * 50)
    print("✅ 视频生成成功！")
    print("=" * 50)
    print(f"📄 视频文件: {output_path}")
    print(f"📊 幻灯片数量: {len(slides_data)} 页")
    print(f"🎤 语音: {voice}")
    print()
    print("💡 提示: 可使用视频播放器查看")
    print("=" * 50)
    
    return {
        'success': True,
        'video_path': output_path,
        'slide_count': len(slides_data),
        'voice': voice
    }


def main():
    """命令行入口"""
    parser = argparse.ArgumentParser(description='视频生成器 - 从文本生成有声PPT视频')
    parser.add_argument('text', nargs='?', help='输入文本内容（或使用--file指定文件）')
    parser.add_argument('--file', '-f', help='从文件读取文本内容')
    parser.add_argument('--pptx', '-p', help='从PPT文件生成视频')
    parser.add_argument('--style', '-s', 
                        choices=['business', 'minimal', 'tech', 'creative'],
                        default='business', help='PPT风格 (default: business)')
    parser.add_argument('--voice', '-v',
                        choices=list(TTS_VOICES.keys()),
                        default='female', help='语音类型 (default: female)')
    parser.add_argument('--output', '-o', help='输出视频路径')
    parser.add_argument('--use-pdf', action='store_true', 
                        help='强制使用PDF中转方案（需要poppler-utils和reportlab）')
    parser.add_argument('--no-pdf', action='store_true',
                        help='强制使用Pillow方案，不使用PDF中转')
    
    args = parser.parse_args()
    
    # 处理PDF中转方案选项
    if args.use_pdf and args.no_pdf:
        print("❌ 错误: --use-pdf 和 --no-pdf 不能同时使用")
        return 1
    
    if args.use_pdf:
        if not check_poppler_available():
            print("❌ 错误: --use-pdf 需要 poppler-utils，请运行: sudo apt-get install poppler-utils")
            return 1
        if not check_reportlab_available():
            print("❌ 错误: --use-pdf 需要 reportlab，请运行: pip install reportlab")
            return 1
        VIDEO_CONFIG['use_pdf_method'] = True
    elif args.no_pdf:
        VIDEO_CONFIG['use_pdf_method'] = False
    
    try:
        if args.pptx:
            result = generate_video_from_pptx(
                pptx_path=args.pptx,
                voice=args.voice,
                output_path=args.output
            )
        else:
            # 获取文本内容
            text = args.text
            if args.file:
                with open(args.file, 'r', encoding='utf-8') as f:
                    text = f.read()
            
            if not text:
                print("❌ 错误: 请提供文本内容或使用 --file 指定文件")
                parser.print_help()
                return 1
            
            result = generate_video_from_text(
                text=text,
                style=args.style,
                voice=args.voice,
                output_path=args.output
            )
        
        return 0 if result['success'] else 1
    
    except Exception as e:
        print(f"❌ 生成失败: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    sys.exit(main())
