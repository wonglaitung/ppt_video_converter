#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成器模块
基于python-pptx库和大模型生成专业PPT演示文稿

风格设计参考业界最佳实践：
- 商务风格：McKinsey/BCG咨询报告风格
- 简约风格：Apple Keynote风格
- 科技风格：NVIDIA GTC/Apple WWDC风格
- 创意风格：Behance/Dribbble热门设计
"""

import os
import sys
import json
import argparse
from datetime import datetime

# 添加项目根目录到路径
script_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(script_dir)
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 导入Qwen大模型接口
try:
    from llm_services.qwen_engine import chat_with_llm
except ImportError:
    print("⚠️ 无法导入Qwen大模型接口，请确保llm_services/qwen_engine.py存在")
    chat_with_llm = None


# ============================================================================
# 业界最佳实践样式配置
# ============================================================================

def get_style_config(style='business'):
    """获取样式配置（参考业界最佳实践）
    
    Args:
        style (str): 模板风格 (business/minimal/tech/creative)
    
    Returns:
        dict: 样式配置字典
    
    风格设计参考：
    - business: McKinsey/BCG咨询报告风格 - 专业稳重、数据驱动
    - minimal: Apple Keynote风格 - 极简美学、大量留白
    - tech: NVIDIA GTC/Apple WWDC风格 - 暗黑科技、霓虹点缀
    - creative: Behance/Dribbble热门设计 - 大胆配色、动态感
    """
    styles = {
        # ====================================================================
        # 商务风格 - McKinsey/BCG咨询报告风格
        # 特点：深蓝渐变、右侧竖条装饰、金色点缀、专业稳重
        # ====================================================================
        'business': {
            'name': '商务风格',
            'description': 'McKinsey/BCG咨询报告风格 - 专业稳重、数据驱动',
            
            # 渐变背景（从左上到右下，淡雅蓝色渐变）
            'gradient': {
                'type': 'linear',
                'angle': 135,
                'stops': [
                    (0.0, RGBColor(240, 245, 250)),    # 极淡蓝
                    (0.5, RGBColor(230, 238, 248)),    # 淡蓝
                    (1.0, RGBColor(215, 228, 242)),    # 浅蓝灰
                ]
            },
            
            # 颜色配置
            'title_color': RGBColor(0, 48, 87),         # 深海蓝（专业）
            'subtitle_color': RGBColor(51, 102, 153),   # 中蓝
            'text_color': RGBColor(51, 51, 51),         # 深灰
            'accent_color': RGBColor(184, 134, 11),     # 金色（点缀）
            'highlight_color': RGBColor(0, 102, 153),   # 蓝色（强调）
            
            # 装饰元素配置
            'decorations': {
                # 右侧竖条装饰（McKinsey经典元素）
                'side_bar': {
                    'enabled': True,
                    'color': RGBColor(0, 82, 147),
                    'width': 0.06,  # 占幻灯片宽度的6%
                    'position': 'right',
                },
                # 底部金色装饰线
                'footer_line': {
                    'enabled': True,
                    'color': RGBColor(184, 134, 11),
                    'height': 0.015,
                },
                # 左上角Logo占位区
                'logo_area': {
                    'enabled': True,
                    'size': 0.5,  # Inches
                },
            },
            
            # 字体配置
            'fonts': {
                'title': {'name': '微软雅黑', 'size': 44, 'bold': True},
                'subtitle': {'name': '微软雅黑', 'size': 24, 'bold': False},
                'heading': {'name': '微软雅黑', 'size': 32, 'bold': True},
                'body': {'name': '微软雅黑', 'size': 20, 'bold': False},
                'caption': {'name': '微软雅黑', 'size': 14, 'bold': False},
            },
        },
        
        # ====================================================================
        # 简约风格 - Apple Keynote风格
        # 特点：纯白背景、极简装饰、大量留白、居中布局
        # ====================================================================
        'minimal': {
            'name': '简约风格',
            'description': 'Apple Keynote风格 - 极简美学、大量留白',
            
            # 纯白背景（Apple经典）
            'gradient': {
                'type': 'solid',
                'color': RGBColor(255, 255, 255),
            },
            
            # 颜色配置
            'title_color': RGBColor(17, 17, 17),         # 纯黑
            'subtitle_color': RGBColor(102, 102, 102),   # 中灰
            'text_color': RGBColor(51, 51, 51),          # 深灰
            'accent_color': RGBColor(0, 122, 255),       # Apple蓝
            'highlight_color': RGBColor(0, 122, 255),    # Apple蓝
            
            # 装饰元素配置（极简）
            'decorations': {
                # 底部细线（唯一的装饰）
                'footer_line': {
                    'enabled': True,
                    'color': RGBColor(230, 230, 230),
                    'height': 0.008,
                },
                # 点缀色块（左下角小方块）
                'accent_block': {
                    'enabled': True,
                    'color': RGBColor(0, 122, 255),
                    'size': 0.15,  # Inches
                    'position': 'bottom_left',
                },
            },
            
            # 字体配置
            'fonts': {
                'title': {'name': '微软雅黑', 'size': 48, 'bold': True},
                'subtitle': {'name': '微软雅黑', 'size': 24, 'bold': False},
                'heading': {'name': '微软雅黑', 'size': 36, 'bold': True},
                'body': {'name': '微软雅黑', 'size': 22, 'bold': False},
                'caption': {'name': '微软雅黑', 'size': 14, 'bold': False},
            },
        },
        
        # ====================================================================
        # 科技风格 - NVIDIA GTC/Apple WWDC风格
        # 特点：深色背景、霓虹青色、网格装饰、未来感
        # ====================================================================
        'tech': {
            'name': '科技风格',
            'description': 'NVIDIA GTC/Apple WWDC风格 - 暗黑科技、霓虹点缀',
            
            # 深色科技渐变背景
            'gradient': {
                'type': 'linear',
                'angle': 180,
                'stops': [
                    (0.0, RGBColor(10, 15, 30)),       # 深蓝黑
                    (0.5, RGBColor(15, 25, 45)),       # 中间蓝
                    (1.0, RGBColor(5, 10, 25)),        # 极深蓝
                ]
            },
            
            # 颜色配置
            'title_color': RGBColor(0, 212, 255),        # 霓虹青
            'subtitle_color': RGBColor(100, 200, 255),   # 浅青
            'text_color': RGBColor(220, 220, 230),       # 浅灰白
            'accent_color': RGBColor(138, 43, 226),      # 紫色点缀
            'highlight_color': RGBColor(0, 255, 200),    # 荧光绿
            
            # 装饰元素配置
            'decorations': {
                # 左侧霓虹边框
                'side_glow': {
                    'enabled': True,
                    'color': RGBColor(0, 212, 255),
                    'width': 0.03,
                    'position': 'left',
                },
                # 右上角网格装饰
                'grid_pattern': {
                    'enabled': True,
                    'color': RGBColor(0, 100, 150),
                    'size': 1.5,  # Inches
                    'opacity': 0.3,
                },
                # 底部渐变线
                'footer_gradient': {
                    'enabled': True,
                    'colors': [RGBColor(0, 212, 255), RGBColor(138, 43, 226)],
                    'height': 0.02,
                },
            },
            
            # 字体配置
            'fonts': {
                'title': {'name': '微软雅黑', 'size': 46, 'bold': True},
                'subtitle': {'name': '微软雅黑', 'size': 24, 'bold': False},
                'heading': {'name': '微软雅黑', 'size': 34, 'bold': True},
                'body': {'name': '微软雅黑', 'size': 20, 'bold': False},
                'caption': {'name': '微软雅黑', 'size': 14, 'bold': False},
            },
        },
        
        # ====================================================================
        # 创意风格 - Behance/Dribbble热门设计
        # 特点：渐变紫粉、几何形状、动态波浪、大胆配色
        # ====================================================================
        'creative': {
            'name': '创意风格',
            'description': 'Behance/Dribbble热门设计 - 大胆配色、动态感',
            
            # 渐变紫粉背景
            'gradient': {
                'type': 'linear',
                'angle': 45,
                'stops': [
                    (0.0, RGBColor(102, 126, 234)),    # 紫蓝
                    (0.5, RGBColor(118, 75, 162)),     # 紫色
                    (1.0, RGBColor(240, 147, 251)),    # 粉紫
                ]
            },
            
            # 颜色配置
            'title_color': RGBColor(255, 255, 255),       # 纯白
            'subtitle_color': RGBColor(255, 240, 245),    # 浅粉白
            'text_color': RGBColor(255, 255, 255),        # 白色
            'accent_color': RGBColor(255, 193, 7),        # 金黄点缀
            'highlight_color': RGBColor(255, 87, 51),     # 橙红强调
            
            # 装饰元素配置
            'decorations': {
                # 右下角圆形装饰
                'corner_circle': {
                    'enabled': True,
                    'color': RGBColor(255, 193, 7),
                    'size': 0.8,  # Inches
                    'position': 'bottom_right',
                },
                # 左上角三角形
                'corner_triangle': {
                    'enabled': True,
                    'color': RGBColor(255, 87, 51),
                    'size': 0.6,  # Inches
                    'position': 'top_left',
                },
                # 底部波浪线（用矩形模拟）
                'wave_accent': {
                    'enabled': True,
                    'color': RGBColor(255, 255, 255),
                    'height': 0.01,
                    'opacity': 0.5,
                },
            },
            
            # 字体配置
            'fonts': {
                'title': {'name': '微软雅黑', 'size': 48, 'bold': True},
                'subtitle': {'name': '微软雅黑', 'size': 26, 'bold': False},
                'heading': {'name': '微软雅黑', 'size': 36, 'bold': True},
                'body': {'name': '微软雅黑', 'size': 22, 'bold': False},
                'caption': {'name': '微软雅黑', 'size': 14, 'bold': False},
            },
        },
    }
    return styles.get(style, styles['business'])


# ============================================================================
# 大模型分析
# ============================================================================

def analyze_text_with_llm(text, enable_thinking=False):
    """使用大模型分析文本生成PPT大纲
    
    Args:
        text (str): 输入文本内容
        enable_thinking (bool): 是否启用推理模式（默认关闭，提升响应速度）
    
    Returns:
        dict: PPT大纲数据
    """
    if not chat_with_llm:
        raise RuntimeError("Qwen大模型接口不可用")
    
    prompt = f"""请分析以下文本内容，生成一个专业的PPT大纲。

【核心原则】
1. **保留关键实体**：要点中必须包含原文中的关键实体名称（人名、公司名、产品名、项目名、地点等）
2. **保留关键数据**：要点中必须包含原文中的关键数据（数字、百分比、金额、日期、指标等）
3. **保持信息完整**：宁可要点长一些，也不要丢失重要信息

【格式要求】
1. 提取一个主标题和副标题，反映文本核心主题
2. 将内容划分为3-8个主要章节，每个章节聚焦一个主题
3. 每个章节包含2-6个要点
4. 每条要点控制在100字以内
5. 只输出JSON格式，不要有任何其他文字

【要点质量示例】

✅ 好的要点（保留了实体和数据）：
- "苹果公司2024年Q4营收949亿美元，同比增长6%，iPhone销量增长5.3%"
- "项目A完成进度85%，预计2024年6月交付，预算执行率92%"
- "北京、上海、广州三地用户占比分别为28%、25%、18%"

❌ 差的要点（丢失了关键信息）：
- "营收增长，销量上升"（缺少具体数据和实体）
- "项目进度良好，按计划推进"（缺少具体进度和日期）

【输出格式】（严格遵守JSON格式）：
{{
  "title": "主标题（简洁有力，5-15字）",
  "subtitle": "副标题（补充说明，10-20字）",
  "sections": [
    {{
      "title": "章节标题（概括主题）",
      "points": [
        "要点1：包含实体名称+具体数据+关键结论",
        "要点2：同上格式",
        "要点3"
      ]
    }},
    {{
      "title": "章节标题2",
      "points": ["要点1", "要点2"]
    }}
  ],
  "conclusion": "核心结论（一句话，包含最重要的结论和关键数据）"
}}

【文本内容】
{text}
"""
    
    response = chat_with_llm(prompt, enable_thinking=enable_thinking)
    
    # 解析JSON响应
    try:
        # 尝试提取JSON部分
        json_str = response.strip()
        if '```json' in json_str:
            json_str = json_str.split('```json')[1].split('```')[0].strip()
        elif '```' in json_str:
            json_str = json_str.split('```')[1].split('```')[0].strip()
        
        outline = json.loads(json_str)
        
        # 验证必需字段
        if 'title' not in outline:
            outline['title'] = '演示文稿'
        if 'sections' not in outline:
            outline['sections'] = []
        
        return outline
    
    except json.JSONDecodeError as e:
        print(f"⚠️ JSON解析失败: {e}")
        print(f"原始响应: {response[:500]}...")
        # 返回默认大纲
        return {
            'title': '演示文稿',
            'subtitle': '自动生成',
            'sections': [
                {'title': '内容概述', 'points': ['请查看原文了解详情']}
            ],
            'conclusion': '感谢观看'
        }


# ============================================================================
# PPT背景和装饰生成函数
# ============================================================================

def set_background(slide, style_config):
    """设置幻灯片背景
    
    Args:
        slide: 幻灯片对象
        style_config (dict): 样式配置
    """
    gradient_config = style_config.get('gradient', {})
    
    background = slide.background
    fill = background.fill
    
    if gradient_config.get('type') == 'solid':
        # 纯色背景
        fill.solid()
        fill.fore_color.rgb = gradient_config['color']
    elif gradient_config.get('type') == 'linear':
        # 渐变背景
        fill.gradient()
        fill.gradient_angle = gradient_config.get('angle', 0)
        
        stops = gradient_config.get('stops', [])
        if len(stops) >= 1:
            fill.gradient_stops[0].position = stops[0][0]
            fill.gradient_stops[0].color.rgb = stops[0][1]
        if len(stops) >= 2:
            fill.gradient_stops[1].position = stops[-1][0]
            fill.gradient_stops[1].color.rgb = stops[-1][1]


def add_shape_with_opacity(slide, shape_type, left, top, width, height, color, opacity=1.0):
    """添加带透明度的形状
    
    Args:
        slide: 幻灯片对象
        shape_type: 形状类型
        left, top, width, height: 位置和尺寸
        color: RGB颜色
        opacity: 透明度 (0-1)
    """
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    
    # 通过brightness模拟透明度效果
    if opacity < 1.0:
        shape.fill.fore_color.brightness = 1.0 - opacity
    
    return shape


def add_decorations(slide, style_config, slide_width, slide_height):
    """添加装饰元素到幻灯片（业界最佳实践）
    
    Args:
        slide: 幻灯片对象
        style_config (dict): 样式配置
        slide_width: 幻灯片宽度（Inches对象）
        slide_height: 幻灯片高度（Inches对象）
    """
    decorations = style_config.get('decorations', {})
    
    # 宽度和高度的数值（英寸）
    w = slide_width.inches
    h = slide_height.inches
    
    # 1. 侧边装饰条（商务/科技风格）
    if 'side_bar' in decorations and decorations['side_bar'].get('enabled'):
        config = decorations['side_bar']
        bar_width = Inches(w * config['width'])
        if config['position'] == 'right':
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                slide_width - bar_width, Inches(0),
                bar_width, slide_height
            )
        else:  # left
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                bar_width, slide_height
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        shape.line.fill.background()
    
    # 2. 侧边霓虹光效（科技风格）
    if 'side_glow' in decorations and decorations['side_glow'].get('enabled'):
        config = decorations['side_glow']
        glow_width = Inches(w * config['width'])
        if config['position'] == 'left':
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                glow_width, slide_height
            )
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                slide_width - glow_width, Inches(0),
                glow_width, slide_height
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        shape.line.fill.background()
    
    # 3. 底部装饰线
    if 'footer_line' in decorations and decorations['footer_line'].get('enabled'):
        config = decorations['footer_line']
        line_height = Inches(h * config['height'])
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), slide_height - line_height,
            slide_width, line_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        shape.line.fill.background()
    
    # 4. 底部渐变线（科技风格）
    if 'footer_gradient' in decorations and decorations['footer_gradient'].get('enabled'):
        config = decorations['footer_gradient']
        line_height = Inches(h * config['height'])
        # 使用第一个颜色作为主色
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), slide_height - line_height,
            slide_width, line_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['colors'][0]
        shape.line.fill.background()
    
    # 5. 底部波浪装饰（创意风格）
    if 'wave_accent' in decorations and decorations['wave_accent'].get('enabled'):
        config = decorations['wave_accent']
        line_height = Inches(h * config['height'])
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), slide_height - line_height - Inches(0.1),
            slide_width, line_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        if config.get('opacity', 1.0) < 1.0:
            shape.fill.fore_color.brightness = 1.0 - config['opacity']
        shape.line.fill.background()
    
    # 6. 角落装饰圆形（创意风格）
    if 'corner_circle' in decorations and decorations['corner_circle'].get('enabled'):
        config = decorations['corner_circle']
        size = Inches(config['size'])
        margin = Inches(0.2)
        if config['position'] == 'bottom_right':
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                slide_width - size - margin, slide_height - size - margin,
                size, size
            )
        elif config['position'] == 'bottom_left':
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                margin, slide_height - size - margin,
                size, size
            )
        elif config['position'] == 'top_right':
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                slide_width - size - margin, margin,
                size, size
            )
        else:  # top_left
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                margin, margin,
                size, size
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        shape.line.fill.background()
    
    # 7. 角落装饰三角形（创意风格）
    if 'corner_triangle' in decorations and decorations['corner_triangle'].get('enabled'):
        config = decorations['corner_triangle']
        size = Inches(config['size'])
        if config['position'] == 'top_left':
            # 用三角形（等腰三角形形状）
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(-size.inches * 0.3), Inches(-size.inches * 0.3),
                size, size
            )
        else:  # bottom_right
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                slide_width - size * 0.7, slide_height - size * 0.7,
                size, size
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        shape.line.fill.background()
    
    # 8. 点缀色块（简约风格）
    if 'accent_block' in decorations and decorations['accent_block'].get('enabled'):
        config = decorations['accent_block']
        size = Inches(config['size'])
        margin = Inches(0.3)
        if config['position'] == 'bottom_left':
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                margin, slide_height - size - margin,
                size, size
            )
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                margin, margin,
                size, size
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        shape.line.fill.background()
    
    # 9. 网格装饰（科技风格）
    if 'grid_pattern' in decorations and decorations['grid_pattern'].get('enabled'):
        config = decorations['grid_pattern']
        size = Inches(config['size'])
        margin = Inches(0.2)
        # 在右上角添加网格
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            slide_width - size - margin, margin,
            size, size
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = config['color']
        if config.get('opacity', 1.0) < 1.0:
            shape.fill.fore_color.brightness = 1.0 - config['opacity']
        shape.line.fill.background()
    
    # 10. Logo占位区域（商务风格）
    if 'logo_area' in decorations and decorations['logo_area'].get('enabled'):
        config = decorations['logo_area']
        size = Inches(config['size'])
        margin = Inches(0.3)
        # 在左上角添加一个小方块作为Logo占位
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            margin, margin,
            size, size * 0.4
        )
        shape.fill.background()  # 透明
        shape.line.color.rgb = style_config.get('accent_color', RGBColor(128, 128, 128))
        shape.line.width = Pt(1)


# ============================================================================
# 幻灯片创建函数
# ============================================================================

def create_title_slide(prs, outline, style_config):
    """创建封面页
    
    Args:
        prs: Presentation对象
        outline (dict): 大纲数据
        style_config (dict): 样式配置
    """
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 设置背景
    set_background(slide, style_config)
    
    # 添加装饰元素
    add_decorations(slide, style_config, slide_width, slide_height)
    
    # 获取字体配置
    fonts = style_config.get('fonts', {})
    title_font = fonts.get('title', {'name': '微软雅黑', 'size': 44, 'bold': True})
    subtitle_font = fonts.get('subtitle', {'name': '微软雅黑', 'size': 24, 'bold': False})
    
    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = outline.get('title', '演示文稿')
    p.font.name = title_font['name']
    p.font.size = Pt(title_font['size'])
    p.font.bold = title_font['bold']
    p.font.color.rgb = style_config['title_color']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if outline.get('subtitle'):
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = outline['subtitle']
        p.font.name = subtitle_font['name']
        p.font.size = Pt(subtitle_font['size'])
        p.font.bold = subtitle_font['bold']
        p.font.color.rgb = style_config.get('subtitle_color', style_config['text_color'])
        p.alignment = PP_ALIGN.CENTER
    
    # 日期
    caption_font = fonts.get('caption', {'name': '微软雅黑', 'size': 14, 'bold': False})
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5)
    )
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = datetime.now().strftime('%Y年%m月%d日')
    p.font.name = caption_font['name']
    p.font.size = Pt(caption_font['size'])
    p.font.bold = caption_font['bold']
    p.font.color.rgb = style_config['accent_color']
    p.alignment = PP_ALIGN.CENTER


def create_toc_slide(prs, outline, style_config):
    """创建目录页
    
    Args:
        prs: Presentation对象
        outline (dict): 大纲数据
        style_config (dict): 样式配置
    """
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 设置背景
    set_background(slide, style_config)
    
    # 添加装饰元素
    add_decorations(slide, style_config, slide_width, slide_height)
    
    # 获取字体配置
    fonts = style_config.get('fonts', {})
    heading_font = fonts.get('heading', {'name': '微软雅黑', 'size': 32, 'bold': True})
    body_font = fonts.get('body', {'name': '微软雅黑', 'size': 20, 'bold': False})
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = '目录'
    p.font.name = heading_font['name']
    p.font.size = Pt(heading_font['size'])
    p.font.bold = heading_font['bold']
    p.font.color.rgb = style_config['title_color']
    
    # 目录内容
    sections = outline.get('sections', [])
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(11.333), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, section in enumerate(sections, 1):
        if i > 1:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]
        
        p.text = f"{i}. {section.get('title', '章节')}"
        p.font.name = body_font['name']
        p.font.size = Pt(body_font['size'])
        p.font.bold = body_font['bold']
        p.font.color.rgb = style_config['text_color']
        p.space_after = Pt(12)
        p.level = 0


def create_content_slide(prs, section, style_config):
    """创建内容页
    
    Args:
        prs: Presentation对象
        section (dict): 章节数据
        style_config (dict): 样式配置
    """
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 设置背景
    set_background(slide, style_config)
    
    # 添加装饰元素
    add_decorations(slide, style_config, slide_width, slide_height)
    
    # 获取字体配置
    fonts = style_config.get('fonts', {})
    heading_font = fonts.get('heading', {'name': '微软雅黑', 'size': 32, 'bold': True})
    body_font = fonts.get('body', {'name': '微软雅黑', 'size': 20, 'bold': False})
    
    # 章节标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = section.get('title', '章节')
    p.font.name = heading_font['name']
    p.font.size = Pt(heading_font['size'])
    p.font.bold = heading_font['bold']
    p.font.color.rgb = style_config['title_color']
    
    # 内容要点
    points = section.get('points', [])
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(11.333), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, point in enumerate(points[:6]):  # 最多6个要点
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]
        
        p.text = f"• {point}"
        p.font.name = body_font['name']
        p.font.size = Pt(body_font['size'])
        p.font.bold = body_font['bold']
        p.font.color.rgb = style_config['text_color']
        p.space_after = Pt(16)
        p.level = 0


def create_conclusion_slide(prs, outline, style_config):
    """创建结尾页
    
    Args:
        prs: Presentation对象
        outline (dict): 大纲数据
        style_config (dict): 样式配置
    """
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 设置背景
    set_background(slide, style_config)
    
    # 添加装饰元素
    add_decorations(slide, style_config, slide_width, slide_height)
    
    # 获取字体配置
    fonts = style_config.get('fonts', {})
    heading_font = fonts.get('heading', {'name': '微软雅黑', 'size': 32, 'bold': True})
    body_font = fonts.get('body', {'name': '微软雅黑', 'size': 20, 'bold': False})
    
    # 总结标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = '总结'
    p.font.name = heading_font['name']
    p.font.size = Pt(heading_font['size'])
    p.font.bold = heading_font['bold']
    p.font.color.rgb = style_config['title_color']
    p.alignment = PP_ALIGN.CENTER
    
    # 总结内容
    conclusion = outline.get('conclusion', '感谢观看')
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(11.333), Inches(1.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    p = content_frame.paragraphs[0]
    p.text = conclusion
    p.font.name = body_font['name']
    p.font.size = Pt(body_font['size'])
    p.font.bold = body_font['bold']
    p.font.color.rgb = style_config['text_color']
    p.alignment = PP_ALIGN.CENTER
    
    # 感谢语
    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(1)
    )
    thanks_frame = thanks_box.text_frame
    p = thanks_frame.paragraphs[0]
    p.text = '感谢观看'
    p.font.name = heading_font['name']
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = style_config['accent_color']
    p.alignment = PP_ALIGN.CENTER


def create_presentation(outline, style='business', output_path=None):
    """根据大纲创建PPT演示文稿
    
    Args:
        outline (dict): PPT大纲数据
        style (str): 模板风格 (business/minimal/tech/creative)
        output_path (str): 输出文件路径
    
    Returns:
        tuple: (ppt文件路径, 大纲json文件路径)
    """
    # 创建演示文稿对象
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 获取样式配置
    style_config = get_style_config(style)
    
    # 1. 创建封面页
    create_title_slide(prs, outline, style_config)
    
    # 2. 创建目录页
    create_toc_slide(prs, outline, style_config)
    
    # 3. 创建内容页
    for section in outline.get('sections', []):
        create_content_slide(prs, section, style_config)
    
    # 4. 创建结尾页
    create_conclusion_slide(prs, outline, style_config)
    
    # 生成输出路径
    if output_path is None:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_dir = os.path.join(project_root, 'output', 'presentations')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'presentation_{timestamp}.pptx')
    
    # 确保输出目录存在
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    
    # 保存PPT文件
    prs.save(output_path)
    
    # 保存大纲JSON文件
    json_path = output_path.replace('.pptx', '.json')
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(outline, f, ensure_ascii=False, indent=2)
    
    return output_path, json_path


# ============================================================================
# 主函数
# ============================================================================

def generate_ppt(text, style='business', output_path=None, enable_thinking=False):
    """生成PPT主函数
    
    Args:
        text (str): 输入文本内容
        style (str): 模板风格
        output_path (str): 输出文件路径
        enable_thinking (bool): 是否启用推理模式（默认关闭，提升响应速度）
    
    Returns:
        dict: 生成结果
    """
    print("=" * 50)
    print("📊 PPT生成器")
    print("=" * 50)
    print(f"📝 输入文本长度: {len(text)} 字符")
    print(f"🎨 模板风格: {get_style_config(style)['name']}")
    print()
    
    # 1. 使用大模型分析文本
    print("🔍 正在分析文本内容...")
    outline = analyze_text_with_llm(text, enable_thinking=enable_thinking)
    
    print(f"📋 提取标题: {outline.get('title', '无')}")
    print(f"📑 提取章节数: {len(outline.get('sections', []))}")
    print()
    
    # 2. 生成PPT文件
    print("🎨 正在生成PPT文件...")
    ppt_path, json_path = create_presentation(outline, style, output_path)
    
    # 3. 统计信息
    slide_count = 2 + len(outline.get('sections', [])) + 1  # 封面+目录+内容+结尾
    
    result = {
        'success': True,
        'ppt_path': ppt_path,
        'json_path': json_path,
        'slide_count': slide_count,
        'outline': outline,
        'style': style
    }
    
    # 4. 输出结果
    print()
    print("=" * 50)
    print("✅ PPT生成成功！")
    print("=" * 50)
    print(f"📄 PPT文件: {ppt_path}")
    print(f"📄 大纲文件: {json_path}")
    print(f"📊 幻灯片数量: {slide_count} 页")
    print(f"📋 标题: {outline.get('title', '无')}")
    print(f"📑 章节数: {len(outline.get('sections', []))}")
    print(f"🎨 风格: {get_style_config(style)['name']}")
    print()
    print("💡 提示: 可使用 PowerPoint 或 WPS 打开查看")
    print("=" * 50)
    
    return result


def main():
    """命令行入口"""
    parser = argparse.ArgumentParser(description='PPT生成器 - 基于大模型生成专业演示文稿')
    parser.add_argument('text', nargs='?', help='输入文本内容（或使用--file指定文件）')
    parser.add_argument('--file', '-f', help='从文件读取文本内容')
    parser.add_argument('--style', '-s', choices=['business', 'minimal', 'tech', 'creative'],
                        default='business', help='模板风格 (default: business)')
    parser.add_argument('--output', '-o', help='输出文件路径')
    parser.add_argument('--thinking', action='store_true', help='启用推理模式（默认关闭以提升响应速度）')
    
    args = parser.parse_args()
    
    # 获取文本内容
    text = args.text
    if args.file:
        with open(args.file, 'r', encoding='utf-8') as f:
            text = f.read()
    
    if not text:
        print("❌ 错误: 请提供文本内容或使用 --file 指定文件")
        parser.print_help()
        return 1
    
    try:
        result = generate_ppt(
            text=text,
            style=args.style,
            output_path=args.output,
            enable_thinking=args.thinking
        )
        return 0 if result['success'] else 1
    
    except Exception as e:
        print(f"❌ 生成失败: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    sys.exit(main())